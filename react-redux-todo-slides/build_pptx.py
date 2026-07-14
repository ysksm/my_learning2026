# -*- coding: utf-8 -*-
"""Build react-redux-todo-slides.pptx from the HTML slide deck design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x17, 0x21, 0x2B)
PAPER = RGBColor(0xFF, 0xFD, 0xF8)
BLUE = RGBColor(0x25, 0x63, 0xEB)
AQUA = RGBColor(0x12, 0xB8, 0xA6)
YELLOW = RGBColor(0xFF, 0xD6, 0x5A)
LINE = RGBColor(0xDF, 0xE7, 0xE8)
CODE_BG = RGBColor(0x12, 0x20, 0x2C)
CODE_FG = RGBColor(0xD7, 0xE7, 0xEF)
CODE_K = RGBColor(0x87, 0xD8, 0xFF)
CODE_C = RGBColor(0x82, 0x98, 0xA8)
BODY = RGBColor(0x48, 0x56, 0x60)
LEAD = RGBColor(0x42, 0x51, 0x5B)
CIRCLE = RGBColor(0xD9, 0xF5, 0xED)
FLOW_BG = RGBColor(0xE8, 0xF6, 0xF5)
TREE_BG = RGBColor(0xF1, 0xF6, 0xF5)
MEM_BG = RGBColor(0xFF, 0xF8, 0xD8)
CHIP_BG = RGBColor(0xDF, 0xF7, 0xF2)
CHIP_FG = RGBColor(0x08, 0x7D, 0x72)
FOOT = RGBColor(0x7A, 0x8A, 0x92)

JP = "Meiryo"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def runs_para(tf, runs, size=14, first=False, align=None, spacing=1.3, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align if align is not None else PP_ALIGN.LEFT
    p.line_spacing = spacing
    p.space_after = Pt(space_after)
    for item in runs:
        text = item[0]
        opts = item[1] if len(item) > 1 else {}
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", False)
        f.name = opts.get("font", JP)
        f.color.rgb = opts.get("color", BODY)
    return p


def box(slide, x, y, w, h, fill=None, line=None, radius=None, shadow_off=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    if shadow_off:
        sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    return sp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb


def new_slide(footer_left, footer_right, circle=True):
    slide = prs.slides.add_slide(BLANK)
    bg = box(slide, -0.05, -0.05, 13.45, 7.6, fill=PAPER)
    if circle:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.4), Inches(-1.8), Inches(3.6), Inches(3.6))
        c.fill.solid()
        c.fill.fore_color.rgb = CIRCLE
        c.line.fill.background()
        c.shadow.inherit = False
    fl = textbox(slide, 0.72, 7.08, 6, 0.3)
    runs_para(fl.text_frame, [(footer_left, {"font": MONO, "color": FOOT})], size=9, first=True, space_after=0)
    fr = textbox(slide, 6.6, 7.08, 6.0, 0.3)
    runs_para(fr.text_frame, [(footer_right, {"font": MONO, "color": FOOT})], size=9, first=True,
              align=PP_ALIGN.RIGHT, space_after=0)
    return slide


def eyebrow(slide, text, y=0.42):
    tb = textbox(slide, 0.72, y, 9, 0.32)
    runs_para(tb.text_frame, [(text, {"bold": True, "color": AQUA})], size=11, first=True, space_after=0)


def heading(slide, text, y=0.78, size=30):
    tb = textbox(slide, 0.72, y, 11.9, 0.75)
    runs_para(tb.text_frame, [(text, {"bold": True, "color": INK})], size=size, first=True, space_after=0)


def lead(slide, runs, y=1.55, w=11.9, size=13):
    tb = textbox(slide, 0.72, y, w, 0.7)
    runs_para(tb.text_frame, [(t, dict({"color": LEAD}, **o)) for t, o in runs], size=size, first=True, space_after=0, spacing=1.4)


def path_chip(slide, text, y=1.42, x=0.72):
    jp = sum(1 for ch in text if ord(ch) > 0x2000)
    w = 0.24 + (len(text) - jp) * 0.092 + jp * 0.17
    sp = box(slide, x, y, w, 0.34, fill=CHIP_BG, radius=0.3)
    sp.text_frame.word_wrap = False
    sp.text_frame.margin_top = Inches(0.02)
    sp.text_frame.margin_bottom = Inches(0.02)
    runs_para(sp.text_frame, [(text, {"font": MONO, "color": CHIP_FG})], size=10, first=True, space_after=0)


def code_block(slide, x, y, w, h, lines, size=11):
    box(slide, x, y, 0.06, h, fill=AQUA)
    sp = box(slide, x, y, w, h, fill=CODE_BG, radius=0.08)
    tf = sp.text_frame
    tf.margin_left = Inches(0.28)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in lines:
        runs = []
        for item in line:
            text, kind = item
            color = {"k": CODE_K, "c": CODE_C}.get(kind, CODE_FG)
            runs.append((text if text else " ", {"font": MONO, "color": color}))
        if not runs:
            runs = [(" ", {"font": MONO, "color": CODE_FG})]
        runs_para(tf, runs, size=size, first=first, space_after=0, spacing=1.25)
        first = False
    return sp


def card(slide, x, y, w, h, blocks):
    """blocks: list of ('h3', text) / ('p', runs) / ('check', [items]) / ('mono', text)"""
    sp = box(slide, x, y, w, h, fill=RGBColor(0xFF, 0xFF, 0xFF), line=LINE, radius=0.07)
    tf = sp.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for kind, content in blocks:
        if kind == "h3":
            runs_para(tf, [(content, {"bold": True, "color": INK})], size=13, first=first, space_after=4)
        elif kind == "p":
            runs_para(tf, content, size=11, first=first, space_after=6, spacing=1.35)
        elif kind == "check":
            for item in content:
                runs_para(tf, [("✓ ", {"bold": True, "color": AQUA})] + item, size=11, first=first, space_after=3, spacing=1.2)
                first = False
            continue
        elif kind == "mono":
            runs_para(tf, [(content, {"font": MONO, "color": CHIP_FG})], size=9.5, first=first, space_after=0)
        first = False
    return sp


def memory(slide, x, y, w, h, runs, size=11):
    box(slide, x, y, 0.07, h, fill=YELLOW)
    sp = box(slide, x + 0.07, y, w - 0.07, h, fill=MEM_BG)
    sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp.text_frame.margin_left = Inches(0.18)
    runs_para(sp.text_frame, runs, size=size, first=True, space_after=0, spacing=1.4)


B = {"bold": True, "color": INK}
CODE_INLINE = {"font": MONO, "color": CHIP_FG}

# ---------------------------------------------------------------- slide 00
s = new_slide("TODO / 00", "START")
tb = textbox(s, 0.72, 1.35, 9, 0.35)
runs_para(tb.text_frame, [("REACT × TYPESCRIPT × REDUX TOOLKIT", {"bold": True, "color": AQUA})], size=12, first=True, space_after=0)
tb = textbox(s, 0.72, 1.8, 8.0, 2.2)
runs_para(tb.text_frame, [("Todo を作って", {"bold": True, "color": INK, "size": 44})], size=44, first=True, space_after=0, spacing=1.15)
runs_para(tb.text_frame, [("Redux を手で覚える。", {"bold": True, "color": BLUE, "size": 44})], size=44, space_after=0, spacing=1.15)
tb = textbox(s, 0.72, 4.15, 7.4, 1.1)
runs_para(tb.text_frame, [("npm create から始め、DDD＋レイヤード構成で「追加・更新・削除」を Redux store に対して実装します。", {"color": LEAD})], size=14, first=True, spacing=1.6, space_after=0)
tag = box(s, 0.72, 5.35, 1.75, 0.42, fill=CHIP_BG, radius=0.5)
tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
runs_para(tag.text_frame, [("所要 60–90 分", {"bold": True, "color": CHIP_FG})], size=11, first=True, align=PP_ALIGN.CENTER, space_after=0)
goal = box(s, 8.5, 3.6, 4.1, 2.9, fill=INK, radius=0.09)
goal.text_frame.margin_left = Inches(0.25)
goal.text_frame.margin_right = Inches(0.25)
goal.text_frame.margin_top = Inches(0.22)
runs_para(goal.text_frame, [("このスライドのゴールは、", {"color": RGBColor(0xFF, 0xFF, 0xFF)})], size=13, first=True, spacing=1.5, space_after=0)
runs_para(goal.text_frame, [("何も見ずに store を組めること。", {"bold": True, "color": YELLOW})], size=13, spacing=1.5, space_after=0)
runs_para(goal.text_frame, [("完成後は Todo の状態が一箇所に集まり、どのコンポーネントからでも安全に操作できます。", {"color": RGBColor(0xFF, 0xFF, 0xFF)})], size=13, spacing=1.5, space_after=0)

# ---------------------------------------------------------------- slide 01
s = new_slide("TODO / 01", "ARCHITECTURE")
eyebrow(s, "01 / 最初に地図を見る")
heading(s, "今回の責務分け")
lead(s, [("小さなアプリでも、状態のルールと画面を分離します。Redux はアプリケーション層の状態管理として扱います。", {})])
flow_items = [("UI", "components"), ("Application", "slice / dispatch"), ("Domain", "Todo 型・ルール"), ("Infrastructure", "Redux store")]
fx = 0.72
fw = 2.62
for i, (t, sub) in enumerate(flow_items):
    fb = box(s, fx, 2.3, fw, 0.95, fill=FLOW_BG, radius=0.18)
    fb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    runs_para(fb.text_frame, [(t, {"bold": True, "color": INK})], size=14, first=True, align=PP_ALIGN.CENTER, space_after=0, spacing=1.1)
    runs_para(fb.text_frame, [(sub, {"color": BODY})], size=10, align=PP_ALIGN.CENTER, space_after=0, spacing=1.1)
    fx += fw
    if i < 3:
        ar = textbox(s, fx + 0.02, 2.45, 0.42, 0.6)
        runs_para(ar.text_frame, [("→", {"bold": True, "color": BLUE, "size": 20})], size=20, first=True, align=PP_ALIGN.CENTER, space_after=0)
        fx += 0.46
card(s, 0.72, 3.55, 3.88, 1.55, [("h3", "Domain"), ("p", [("Todo が何であるか。型や純粋なルールを置く。", {})])])
card(s, 4.73, 3.55, 3.88, 1.55, [("h3", "Application"), ("p", [("画面が使う操作をまとめる。今回は slice が中心。", {})])])
card(s, 8.74, 3.55, 3.88, 1.55, [("h3", "Presentation"), ("p", [("表示とイベント受付だけ。Redux の詳細を持ち込まない。", {})])])
memory(s, 0.72, 5.5, 11.9, 0.85, [("覚える核：", B), ("state → slice（操作を定義）→ store（登録）→ Provider（配る）→ hooks（読む・送る）", {"color": LEAD})])

# ---------------------------------------------------------------- slide 02
s = new_slide("TODO / 02", "SETUP")
eyebrow(s, "02 / プロジェクトを作る")
heading(s, "Vite で React + TypeScript を開始")
code_block(s, 0.72, 1.7, 6.6, 2.15, [
    [("# プロジェクト作成", "c")],
    [("npm create vite@latest todo-redux -- --template react-ts", None)],
    [("cd todo-redux", None)],
    [("npm install", None)],
    [("npm install @reduxjs/toolkit react-redux", None)],
    [("npm run dev", None)],
], size=11)
memory(s, 0.72, 4.1, 6.6, 1.15, [
    ("@reduxjs/toolkit", B), (" は slice と store を作る。\n", {"color": LEAD}),
    ("react-redux", B), (" は React と Redux をつなぐ。", {"color": LEAD})])
card(s, 7.6, 1.7, 5.0, 3.9, [
    ("h3", "作成時の選択肢"),
    ("check", [[("Framework: React", {"color": LEAD})],
               [("Variant: TypeScript + SWC でOK", {"color": LEAD})],
               [("ブラウザで Vite の初期画面が見えたら成功", {"color": LEAD})]]),
    ("h3", "最初の確認"),
    ("p", [("src/main.tsx", CODE_INLINE), (" がアプリの入口です。ここで後ほど ", {}), ("Provider", CODE_INLINE), (" を追加します。", {})]),
])

# ---------------------------------------------------------------- slide 03
s = new_slide("TODO / 03", "FOLDERS")
eyebrow(s, "03 / 置き場所を決める")
heading(s, "機能単位 × レイヤーで整理する")
tree = box(s, 0.72, 1.7, 6.6, 3.6, fill=TREE_BG, radius=0.06)
tf = tree.text_frame
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.2)
tree_lines = [
    "src/",
    "├── app/",
    "│   └── store.ts             # Redux の組み立て",
    "├── features/todo/",
    "│   ├── domain/todo.ts       # Todo の型",
    "│   ├── application/todoSlice.ts # 操作と状態",
    "│   └── presentation/TodoApp.tsx # 画面",
    "├── App.tsx",
    "└── main.tsx",
]
first = True
for ln in tree_lines:
    runs_para(tf, [(ln, {"font": MONO, "color": RGBColor(0x25, 0x39, 0x47)})], size=11, first=first, space_after=0, spacing=1.35)
    first = False
card(s, 7.6, 1.7, 5.0, 3.6, [
    ("h3", "この分け方の意図"),
    ("p", [("features/todo", B), (" に Todo 関連を閉じ込めます。別の機能（auth など）が増えても迷いません。", {})]),
    ("p", [("本格的な DDD では repository や use case を増やせます。今回は Redux を覚えるため、最小構成に絞ります。", {})]),
    ("mono", "mkdir -p src/app src/features/todo/{domain,application,presentation}"),
])

# ---------------------------------------------------------------- slide 04
s = new_slide("TODO / 04", "DOMAIN")
eyebrow(s, "04 / Domain")
heading(s, "まず Todo を言葉にする")
path_chip(s, "src/features/todo/domain/todo.ts", y=1.5)
code_block(s, 0.72, 2.0, 8.2, 2.4, [
    [("export type", "k"), (" Todo = {", None)],
    [("  id: string;", None)],
    [("  title: string;", None)],
    [("  completed: boolean;", None)],
    [("};", None)],
    [("", None)],
    [("export type", "k"), (" CreateTodoInput = { title: string };", None)],
], size=11)
card(s, 0.72, 4.7, 5.85, 1.75, [
    ("h3", "なぜ ID が必要？"),
    ("p", [("更新・削除する対象を一意に指定するため。配列の index は並び替えで変わるので使いません。", {})])])
card(s, 6.77, 4.7, 5.85, 1.75, [
    ("h3", "なぜ input 型を分ける？"),
    ("p", [("新規作成時には id と completed を画面から渡させない、というルールを型で示せます。", {})])])

# ---------------------------------------------------------------- slide 05
s = new_slide("TODO / 05", "SLICE")
eyebrow(s, "05 / Slice の骨組み")
heading(s, "state と操作を一緒に定義する")
path_chip(s, "src/features/todo/application/todoSlice.ts", y=1.5)
code_block(s, 0.72, 2.0, 11.2, 3.35, [
    [("import", "k"), (" { createSlice, ", None), ("type", "k"), (" PayloadAction } ", None), ("from", "k"), (" '@reduxjs/toolkit';", None)],
    [("import type", "k"), (" { Todo, CreateTodoInput } ", None), ("from", "k"), (" '../domain/todo';", None)],
    [("", None)],
    [("type", "k"), (" TodoState = { items: Todo[] };", None)],
    [("const", "k"), (" initialState: TodoState = { items: [] };", None)],
    [("", None)],
    [("const", "k"), (" todoSlice = createSlice({", None)],
    [("  name: 'todo', initialState,", None)],
    [("  reducers: { ", None), ("/* 次のスライドで CRUD を追加 */", "c"), (" },", None)],
    [("});", None)],
], size=11)
memory(s, 0.72, 5.65, 11.9, 0.95, [
    ("createSlice が作るもの：", B),
    ("reducer と action creator。状態・操作名・状態を変える処理を同じファイルに置くのが RTK の基本です。", {"color": LEAD})])

# ---------------------------------------------------------------- slide 06
s = new_slide("TODO / 06", "CRUD")
eyebrow(s, "06 / CRUD reducers")
heading(s, "追加・更新・削除を定義する")
path_chip(s, "todoSlice.ts の reducers 内", y=1.5)
code_block(s, 0.72, 1.95, 11.9, 3.85, [
    [("addTodo: (state, action: PayloadAction<CreateTodoInput>) => {", None)],
    [("  state.items.push({ id: crypto.randomUUID(), title: action.payload.title, completed: false });", None)],
    [("},", None)],
    [("toggleTodo: (state, action: PayloadAction<string>) => {", None)],
    [("  ", None), ("const", "k"), (" todo = state.items.find(item => item.id === action.payload);", None)],
    [("  ", None), ("if", "k"), (" (todo) todo.completed = !todo.completed;", None)],
    [("},", None)],
    [("renameTodo: (state, action: PayloadAction<{ id: string; title: string }>) => {", None)],
    [("  ", None), ("const", "k"), (" todo = state.items.find(item => item.id === action.payload.id);", None)],
    [("  ", None), ("if", "k"), (" (todo) todo.title = action.payload.title;", None)],
    [("},", None)],
    [("deleteTodo: (state, action: PayloadAction<string>) => {", None)],
    [("  state.items = state.items.filter(item => item.id !== action.payload);", None)],
    [("},", None)],
], size=10)
memory(s, 0.72, 6.0, 11.9, 0.9, [
    ("Immer：", B),
    ("RTK が安全な不変更新に変換するため、push やプロパティ代入で書けます。更新は「id + 新しい値」、削除は「id」だけを payload に渡します。", {"color": LEAD})], size=10.5)

# ---------------------------------------------------------------- slide 07
s = new_slide("TODO / 07", "EXPORT")
eyebrow(s, "07 / Slice を完成させる")
heading(s, "action と reducer を export する")
path_chip(s, "todoSlice.ts の末尾", y=1.5)
code_block(s, 0.72, 2.0, 9.4, 1.75, [
    [("export const", "k"), (" { addTodo, toggleTodo, renameTodo, deleteTodo } =", None)],
    [("  todoSlice.actions;", None)],
    [("", None)],
    [("export default", "k"), (" todoSlice.reducer;", None)],
], size=12)
card(s, 0.72, 4.2, 5.85, 1.85, [
    ("h3", "名前付き export"),
    ("p", [("コンポーネントが ", {}), ("dispatch(addTodo(...))", CODE_INLINE), (" できる action creator。", {})])])
card(s, 6.77, 4.2, 5.85, 1.85, [
    ("h3", "default export"),
    ("p", [("store に登録する reducer。slice ごとに一つです。", {})])])

# ---------------------------------------------------------------- slide 08
s = new_slide("TODO / 08", "STORE")
eyebrow(s, "08 / Store")
heading(s, "slice を一つの store に登録する")
path_chip(s, "src/app/store.ts", y=1.5)
code_block(s, 0.72, 2.0, 10.6, 3.3, [
    [("import", "k"), (" { configureStore } ", None), ("from", "k"), (" '@reduxjs/toolkit';", None)],
    [("import", "k"), (" todoReducer ", None), ("from", "k"), (" '../features/todo/application/todoSlice';", None)],
    [("", None)],
    [("export const", "k"), (" store = configureStore({", None)],
    [("  reducer: { todo: todoReducer },", None)],
    [("});", None)],
    [("", None)],
    [("export type", "k"), (" RootState = ReturnType<", None), ("typeof", "k"), (" store.getState>;", None)],
    [("export type", "k"), (" AppDispatch = ", None), ("typeof", "k"), (" store.dispatch;", None)],
], size=11)
memory(s, 0.72, 5.6, 11.9, 0.9, [
    ("暗記ポイント：", B),
    ("reducer: { todo: todoReducer }", CODE_INLINE),
    (" のキーが、", {"color": LEAD}),
    ("state.todo", CODE_INLINE),
    (" の名前になります。", {"color": LEAD})])

# ---------------------------------------------------------------- slide 09
s = new_slide("TODO / 09", "PROVIDER")
eyebrow(s, "09 / Provider")
heading(s, "React 全体に store を配る")
path_chip(s, "src/main.tsx", y=1.5)
code_block(s, 0.72, 2.0, 10.6, 2.75, [
    [("import", "k"), (" { Provider } ", None), ("from", "k"), (" 'react-redux';", None)],
    [("import", "k"), (" { store } ", None), ("from", "k"), (" './app/store';", None)],
    [("import", "k"), (" App ", None), ("from", "k"), (" './App';", None)],
    [("", None)],
    [("createRoot(document.getElementById('root')!).render(", None)],
    [("  <StrictMode><Provider store={store}><App /></Provider></StrictMode>,", None)],
    [(");", None)],
], size=11)
memory(s, 0.72, 5.1, 11.9, 0.9, [
    ("Provider がないと", B),
    (" useSelector と useDispatch は使えません。アプリの入口で一度だけ囲みます。", {"color": LEAD})])

# ---------------------------------------------------------------- slide 10
s = new_slide("TODO / 10", "READ / DISPATCH")
eyebrow(s, "10 / Presentation")
heading(s, "画面から state を読み、action を送る")
path_chip(s, "src/features/todo/presentation/TodoApp.tsx", y=1.5)
code_block(s, 0.72, 1.95, 11.9, 4.0, [
    [("const", "k"), (" [title, setTitle] = useState('');", None)],
    [("const", "k"), (" todos = useSelector((state: RootState) => state.todo.items);", None)],
    [("const", "k"), (" dispatch = useDispatch<AppDispatch>();", None)],
    [("", None)],
    [("<form onSubmit={(e) => {", None)],
    [("  e.preventDefault();", None)],
    [("  ", None), ("if", "k"), (" (title.trim()) { dispatch(addTodo({ title: title.trim() })); setTitle(''); }", None)],
    [("}}>", None)],
    [("  <input value={title} onChange={e => setTitle(e.target.value)} />", None)],
    [("  <button>追加</button>", None)],
    [("</form>", None)],
    [("{todos.map(todo => <li key={todo.id}>", None)],
    [("  <input type=\"checkbox\" checked={todo.completed} onChange={() => dispatch(toggleTodo(todo.id))} />", None)],
    [("  {todo.title} <button onClick={() => dispatch(deleteTodo(todo.id))}>削除</button>", None)],
    [("</li>)}", None)],
], size=10)
memory(s, 0.72, 6.15, 11.9, 0.75, [
    ("流れ：", B),
    ("イベント → dispatch(action(payload)) → reducer → store 更新 → useSelector の画面が再描画。", {"color": LEAD})], size=10.5)

# ---------------------------------------------------------------- slide 11
s = new_slide("TODO / 11", "YOU CAN BUILD IT")
eyebrow(s, "11 / 仕上げ")
heading(s, "見ずに再現するためのチェック")
card(s, 0.72, 1.7, 5.85, 3.85, [
    ("h3", "15 分チャレンジ"),
    ("check", [
        [("Todo 型を書ける", {"color": LEAD})],
        [("createSlice", CODE_INLINE), (" の3要素を書ける", {"color": LEAD})],
        [("4 action を export できる", {"color": LEAD})],
        [("configureStore", CODE_INLINE), (" に reducer を登録できる", {"color": LEAD})],
        [("Provider", CODE_INLINE), (" でアプリを囲める", {"color": LEAD})],
        [("useSelector", CODE_INLINE), (" / ", {"color": LEAD}), ("useDispatch", CODE_INLINE), (" で CRUD できる", {"color": LEAD})],
    ]),
])
card(s, 6.77, 1.7, 5.85, 3.85, [
    ("h3", "つまずきやすい点"),
    ("p", [("state.todo が undefined", B), ("\nstore のキーと selector が一致しているか確認。", {})]),
    ("p", [("dispatch が使えない", B), ("\nmain.tsx の Provider を確認。", {})]),
    ("p", [("更新されない", B), ("\nid を payload に渡しているか、対象 Todo を find できているか確認。", {})]),
])
memory(s, 0.72, 5.85, 11.9, 0.9, [
    ("次の一歩：", B),
    ("編集フォームで ", {"color": LEAD}),
    ("renameTodo({ id, title })", CODE_INLINE),
    (" を dispatch し、localStorage 永続化や async thunk に広げましょう。", {"color": LEAD})])

out = "/home/user/my_learning2026/react-redux-todo-slides/react-redux-todo-slides.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
