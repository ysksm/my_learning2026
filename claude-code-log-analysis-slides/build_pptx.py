# -*- coding: utf-8 -*-
"""Build claude-code-log-analysis-slides.pptx (matches index.html deck)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x17, 0x21, 0x2B)
PAPER = RGBColor(0xFF, 0xFD, 0xF8)
BLUE = RGBColor(0x25, 0x63, 0xEB)
AQUA = RGBColor(0x12, 0xB8, 0xA6)
YELLOW = RGBColor(0xFF, 0xD6, 0x5A)
LINE = RGBColor(0xDF, 0xE7, 0xE8)
CODE_BG = RGBColor(0x12, 0x20, 0x2C)
CODE_FG = RGBColor(0xD7, 0xE7, 0xEF)
CODE_K = RGBColor(0x87, 0xD8, 0xFF)
CODE_C = RGBColor(0x82, 0x98, 0xA8)
BODY = RGBColor(0x48, 0x56, 0x60)
LEAD = RGBColor(0x42, 0x51, 0x5B)
CIRCLE = RGBColor(0xD9, 0xF5, 0xED)
FLOW_BG = RGBColor(0xE8, 0xF6, 0xF5)
MEM_BG = RGBColor(0xFF, 0xF8, 0xD8)
CHIP_FG = RGBColor(0x08, 0x7D, 0x72)
FOOT = RGBColor(0x7A, 0x8A, 0x92)
TH_BG = RGBColor(0xE8, 0xF6, 0xF5)

JP = "Meiryo"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def base_slide(footer_left, footer_right):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background()
    bg.shadow.inherit = False
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(-2.0), Inches(3.6), Inches(3.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = CIRCLE; circ.line.fill.background()
    circ.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.6), Inches(7.06), Inches(12.1), Inches(0.35))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = footer_left
    r.font.size = Pt(10); r.font.name = MONO; r.font.color.rgb = FOOT
    p2 = tf.paragraphs[0]
    tb2 = s.shapes.add_textbox(Inches(8.7), Inches(7.06), Inches(4.0), Inches(0.35))
    p3 = tb2.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.RIGHT
    r = p3.add_run(); r.text = footer_right
    r.font.size = Pt(10); r.font.name = MONO; r.font.color.rgb = FOOT
    return s


def text(s, x, y, w, h, runs, size=14, bold=False, color=BODY, font=JP,
         align=PP_ALIGN.LEFT, spacing=1.25):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        items = line if isinstance(line, list) else [(line, {})]
        for t, opts in [(i[0], i[1] if len(i) > 1 else {}) for i in items]:
            r = p.add_run(); r.text = t
            f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.name = opts.get("font", font)
            f.color.rgb = opts.get("color", color)
    return tb


def eyebrow(s, label):
    text(s, 0.6, 0.32, 11, 0.35, label, size=12, bold=True, color=AQUA, font=MONO)


def h2(s, title, y=0.68):
    text(s, 0.6, y, 12.0, 0.75, title, size=30, bold=True, color=INK)


def lead(s, body, y=1.42, w=12.0, size=13):
    text(s, 0.6, y, w, 0.8, body, size=size, color=LEAD)


def codebox(s, x, y, w, h, lines, size=11):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG; box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.12)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.18
        items = line if isinstance(line, list) else [(line, {})]
        for t, opts in [(i[0], i[1] if len(i) > 1 else {}) for i in items]:
            r = p.add_run(); r.text = t
            f = r.font
            f.size = Pt(opts.get("size", size)); f.name = MONO
            f.color.rgb = opts.get("color", CODE_FG)
    return box


def card(s, x, y, w, h, title, bullets, tsize=14, bsize=12):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = True; r.font.name = JP; r.font.color.rgb = INK
    for b in bullets:
        p = tf.add_paragraph(); p.line_spacing = 1.3; p.space_before = Pt(4)
        mark, tcol = ("✓ ", AQUA) if not b.startswith("!") else ("", BODY)
        b = b.lstrip("!")
        r = p.add_run(); r.text = mark; r.font.size = Pt(bsize); r.font.bold = True
        r.font.color.rgb = tcol; r.font.name = JP
        r = p.add_run(); r.text = b
        r.font.size = Pt(bsize); r.font.name = JP; r.font.color.rgb = BODY
    return box


def memory(s, x, y, w, h, label, body, size=12):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    bar.shadow.inherit = False
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.07), Inches(y), Inches(w - 0.07), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = MEM_BG; box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.09)
    p = tf.paragraphs[0]; p.line_spacing = 1.3
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.name = JP; r.font.color.rgb = INK
    r = p.add_run(); r.text = body
    r.font.size = Pt(size); r.font.name = JP; r.font.color.rgb = RGBColor(0x4A, 0x44, 0x20)


def table(s, x, y, w, col_widths, rows, hsize=12, bsize=11, row_h=0.34):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = s.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
    tbl = shape.table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH_BG if ri == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]; p.line_spacing = 1.1
            r = p.add_run(); r.text = cell_text
            f = r.font
            f.size = Pt(hsize if ri == 0 else bsize)
            f.bold = ri == 0
            f.name = JP
            f.color.rgb = INK if ri == 0 else BODY
    return shape


def flow(s, x, y, w, h, steps):
    n = len(steps)
    gap = 0.42
    bw = (w - gap * (n - 1)) / n
    for i, (t1, t2) in enumerate(steps):
        bx = x + i * (bw + gap)
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(bw), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = FLOW_BG; b.line.fill.background()
        b.shadow.inherit = False
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t1
        r.font.size = Pt(13); r.font.bold = True; r.font.name = JP; r.font.color.rgb = INK
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t2
        r.font.size = Pt(10); r.font.name = MONO; r.font.color.rgb = RGBColor(0x4C, 0x6A, 0x68)
        if i < n - 1:
            text(s, bx + bw + 0.04, y + h / 2 - 0.22, gap, 0.4, "→",
                 size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 00 TITLE
s = base_slide("CC-LOG / 00", "START")
text(s, 0.6, 0.9, 8.3, 0.4, "CLAUDE CODE × JSONL × DATA ANALYSIS",
     size=13, bold=True, color=AQUA, font=MONO)
text(s, 0.6, 1.35, 8.6, 2.2, [
    [("Claude Code のログを", {"size": 42, "bold": True, "color": INK})],
    [("丸ごと解析する。", {"size": 42, "bold": True, "color": BLUE})],
])
text(s, 0.6, 3.4, 7.6, 1.9,
     "Claude Code は全セッションを JSONL で記録しています。その在り処・スキーマ・分析手順を "
     "step by step で辿り、会話・ツール・トークン・時間・エラーまで「解析できるもの全部」を取り出します。",
     size=14, color=LEAD)
chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.15), Inches(2.5), Inches(0.42))
chip.fill.solid(); chip.fill.fore_color.rgb = RGBColor(0xE0, 0xF6, 0xF2); chip.line.fill.background()
chip.shadow.inherit = False
p = chip.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "jq + Python だけでOK"
r.font.size = Pt(12); r.font.bold = True; r.font.name = JP; r.font.color.rgb = CHIP_FG
goal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(2.6), Inches(4.15), Inches(3.4))
goal.fill.solid(); goal.fill.fore_color.rgb = INK; goal.line.fill.background()
goal.shadow.inherit = False
tf = goal.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.22); tf.margin_right = Inches(0.22)
p = tf.paragraphs[0]; p.line_spacing = 1.5
r = p.add_run(); r.text = "このスライドのゴールは、"
r.font.size = Pt(13); r.font.name = JP; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p = tf.add_paragraph(); p.line_spacing = 1.5
r = p.add_run(); r.text = "自分のログから利用実態レポートを自力で作れること。"
r.font.size = Pt(14); r.font.bold = True; r.font.name = JP; r.font.color.rgb = YELLOW
p = tf.add_paragraph(); p.line_spacing = 1.5
r = p.add_run(); r.text = "よく使うツール、かかったコスト、待ち時間、失敗傾向まで数字で見えるようになります。"
r.font.size = Pt(13); r.font.name = JP; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------- 01 WHERE
s = base_slide("CC-LOG / 01", "WHERE")
eyebrow(s, "01 / まず在り処を知る")
h2(s, "ログはどこにある？ — ~/.claude/ の地図")
lead(s, "主役はプロジェクトごとのセッショントランスクリプト。作業ディレクトリのパスを / → - に変換したフォルダに、セッションIDごとの JSONL が置かれます。")
codebox(s, 0.6, 2.15, 12.1, 3.5, [
    "~/.claude/",
    [("├── projects/", {}), ("                       ← ★ 本命：セッショントランスクリプト", {"color": CODE_K})],
    "│   └── -home-user-my_learning2026/   (cwd の / を - に置換した名前)",
    "│       └── <session-id>.jsonl        (1セッション = 1ファイル)",
    [("├── history.jsonl", {}), ("                   ← 入力プロンプトの履歴", {"color": CODE_K})],
    [("├── todos/", {}), ("                          ← TodoWrite のタスク状態 (JSON)", {"color": CODE_K})],
    [("├── shell-snapshots/", {}), ("                ← Bash 実行時のシェル環境スナップショット", {"color": CODE_K})],
    [("└── settings.json", {}), ("                   ← cleanupPeriodDays で保持期間を制御", {"color": CODE_K})],
], size=12)
memory(s, 0.6, 5.95, 12.1, 0.75, "注意：",
       "古いトランスクリプトは cleanupPeriodDays（既定30日）で自動削除されます。長期分析したい場合は先にコピーして退避しましょう。")

# ---------------------------------------------------------------- 02 STEP1
s = base_slide("CC-LOG / 02", "STEP 1 — INVENTORY")
eyebrow(s, "02 / STEP 1")
h2(s, "JSONL の基本 — 1行 = 1レコード")
lead(s, "JSONL は「1行に1つの JSON オブジェクト」を並べた形式。行単位で読めるので jq やストリーム処理と相性抜群。まずはセッション一覧と行数から。")
codebox(s, 0.6, 2.2, 12.1, 2.3, [
    [("# プロジェクトのログディレクトリへ", {"color": CODE_C})],
    "cd ~/.claude/projects/-home-user-my_learning2026",
    [("# セッション一覧（更新順）とレコード数", {"color": CODE_C})],
    "ls -lt *.jsonl",
    "wc -l *.jsonl",
    [("# 最初の1レコードを整形して眺める", {"color": CODE_C})],
    "head -1 <session-id>.jsonl | jq .",
])
card(s, 0.6, 4.75, 5.9, 2.0, "ここで分かること", [
    "セッションの数と規模（行数 ≒ イベント数）",
    "ファイル更新日時 ＝ 最後に使った日",
    "プロジェクトごとの利用量の偏り",
])
card(s, 6.8, 4.75, 5.9, 2.0, "準備するツール", [
    "!jq（集計ワンライナー用）",
    "!Python 3（本格分析用。pandas があると便利）",
])

# ---------------------------------------------------------------- 03 STEP2
s = base_slide("CC-LOG / 03", "STEP 2 — RECORD TYPES")
eyebrow(s, "03 / STEP 2")
h2(s, "レコード種別を把握する — type フィールド")
lead(s, "全レコードはトップレベルの type で分類できます。まず種別の分布を数えると、そのセッションの構造が一目で分かります。")
codebox(s, 0.6, 2.1, 12.1, 0.55, ["jq -r '.type' *.jsonl | sort | uniq -c | sort -rn"])
table(s, 0.6, 2.95, 12.1, [2.3, 3.4, 6.4], [
    ["type", "意味", "中身のポイント"],
    ["user", "ユーザー発話 or ツール結果", "content が文字列なら人間の入力、配列（tool_result）ならツールの戻り"],
    ["assistant", "Claude の応答", "content[] に text / thinking / tool_use。usage（トークン）と model 付き"],
    ["attachment", "コンテキスト添付", "システムリマインダーや差し込み情報"],
    ["queue-operation", "入力キュー操作", "enqueue / dequeue とプロンプト原文"],
    ["summary / last-prompt 等", "セッション管理用", "コンパクト時の要約、再開用ポインタ"],
], row_h=0.62)

# ---------------------------------------------------------------- 04 STEP3
s = base_slide("CC-LOG / 04", "STEP 3 — METADATA")
eyebrow(s, "04 / STEP 3")
h2(s, "共通メタデータを読む")
lead(s, "user / assistant レコードには毎回同じメタデータが付きます。これが後の全分析の「軸」になります。")
table(s, 0.6, 2.15, 6.1, [2.5, 3.6], [
    ["フィールド", "使い道"],
    ["uuid / parentUuid", "会話ツリーの復元（次ステップ）"],
    ["timestamp", "時間分析すべての基礎（ISO 8601）"],
    ["sessionId", "セッション横断の集計キー"],
    ["cwd / gitBranch", "どのリポジトリ・ブランチの作業か"],
    ["version", "Claude Code のバージョン推移"],
    ["isSidechain", "サブエージェントの発話か"],
    ["permissionMode", "権限モード（default / plan 等）"],
    ["requestId", "API リクエスト単位の紐付け"],
], row_h=0.47)
codebox(s, 7.0, 2.15, 5.7, 3.6, [
    [("# 誰が・いつ・どこで", {"color": CODE_C})],
    "{",
    [('  "type": ', {}), ('"user",', {"color": CODE_K})],
    '  "uuid": "81ac19aa-…",',
    '  "parentUuid": null,',
    '  "timestamp": "2026-07-15T11:56:36Z",',
    '  "sessionId": "8c908057-…",',
    '  "cwd": "/home/user/my_learning2026",',
    '  "gitBranch": "claude/…-x3180r",',
    '  "version": "2.1.210"',
    "}",
], size=11)
memory(s, 0.6, 6.35, 12.1, 0.6, "ポイント：",
       "gitBranch と cwd があるので「どのブランチ作業に何トークン使ったか」まで割り出せます。")

# ---------------------------------------------------------------- 05 STEP4
s = base_slide("CC-LOG / 05", "STEP 4 — CONVERSATION TREE")
eyebrow(s, "05 / STEP 4")
h2(s, "会話ツリーを復元する — parentUuid")
lead(s, "各レコードは parentUuid で直前のレコードを指す連結リスト（分岐すると木）。これを辿ると会話の実際の流れが再現できます。")
flow(s, 0.6, 2.2, 12.1, 0.95, [
    ("user", "prompt (parent=null)"),
    ("assistant", "thinking / text"),
    ("assistant", "tool_use"),
    ("user", "tool_result"),
    ("assistant", "text = 最終回答"),
])
codebox(s, 0.6, 3.45, 12.1, 2.3, [
    [("import", {"color": CODE_K}), (" json", {})],
    "nodes = {}",
    [("with", {"color": CODE_K}), (" open(path) ", {}), ("as", {"color": CODE_K}), (" f:", {})],
    [("    for", {"color": CODE_K}), (" line ", {}), ("in", {"color": CODE_K}), (" f:", {})],
    "        d = json.loads(line)",
    [("        if", {"color": CODE_K}), (" 'uuid' ", {}), ("in", {"color": CODE_K}), (" d: nodes[d['uuid']] = d", {})],
    [("def", {"color": CODE_K}), (" chain(leaf):  ", {}), ("# 葉から根まで遡る", {"color": CODE_C})],
    [("    while", {"color": CODE_K}), (" leaf: ", {}), ("yield", {"color": CODE_K}),
     (" leaf; leaf = nodes.get(leaf.get('parentUuid') ", {}), ("or", {"color": CODE_K}), (" '')", {})],
])
memory(s, 0.6, 6.05, 12.1, 0.75, "sidechain：",
       "isSidechain: true の系列はサブエージェント（Task/Agent）の独立した会話。本流と分けて集計すると「委譲した仕事量」が測れます。")

# ---------------------------------------------------------------- 06 STEP5
s = base_slide("CC-LOG / 06", "STEP 5 — PROMPTS")
eyebrow(s, "06 / STEP 5")
h2(s, "ユーザープロンプトを抽出する")
lead(s, "「人間が何を頼んだか」は分析の出発点。type=user のうち content が文字列のものが人間の入力、配列のものはツール結果なので除外します。")
codebox(s, 0.6, 2.2, 12.1, 1.75, [
    [("# 人間のプロンプトだけを時刻付きで抽出", {"color": CODE_C})],
    'jq -r \'select(.type=="user" and (.message.content|type)=="string")',
    "       | [.timestamp, .message.content] | @tsv' *.jsonl",
    [("# 全プロジェクト横断ならこちら（入力履歴）", {"color": CODE_C})],
    "jq -r '.display' ~/.claude/history.jsonl | head -50",
])
card(s, 0.6, 4.25, 5.9, 2.25, "ここで分かること", [
    "依頼の種類（実装 / 調査 / 修正 / 質問）の割合",
    "1セッションあたりの指示回数 ＝ 往復数",
    "プロンプトの長さと成果の関係",
])
card(s, 6.8, 4.25, 5.9, 2.25, "ワンランク上", [
    "!抽出したプロンプトを Claude に分類させる（「このプロンプト集を用途別にタグ付けして」）と、自分の使い方の癖が定量化できます。",
])

# ---------------------------------------------------------------- 07 STEP6
s = base_slide("CC-LOG / 07", "STEP 6 — TOOLS")
eyebrow(s, "07 / STEP 6")
h2(s, "ツール使用を分析する — tool_use と tool_result")
lead(s, "assistant の tool_use（name / input / id）と、user 側の tool_result（同じ tool_use_id）をペアにすると Claude の「行動ログ」になります。")
codebox(s, 0.6, 2.2, 12.1, 2.15, [
    [("# ツール別の使用回数ランキング", {"color": CODE_C})],
    'jq -r \'select(.type=="assistant") | .message.content[]?',
    '       | select(.type=="tool_use") | .name\' *.jsonl | sort | uniq -c | sort -rn',
    [("# Bash で実行されたコマンドの傾向（先頭単語）", {"color": CODE_C})],
    'jq -r \'… select(.type=="tool_use" and .name=="Bash") | .input.command\' *.jsonl',
    "  | awk '{print $1}' | sort | uniq -c | sort -rn",
])
card(s, 0.6, 4.65, 3.9, 1.95, "行動の内訳", [
    "!Read / Edit / Write / Bash / Grep… の比率で「読む仕事か書く仕事か」が見える。",
])
card(s, 4.7, 4.65, 3.9, 1.95, "触ったファイル", [
    "!input.file_path を集計すればホットスポットのファイルが特定できる。",
])
card(s, 8.8, 4.65, 3.9, 1.95, "実行結果", [
    "!toolUseResult の stdout / stderr / interrupted で成否と中断を判定。",
])

# ---------------------------------------------------------------- 08 STEP7
s = base_slide("CC-LOG / 08", "STEP 7 — TOKENS & COST")
eyebrow(s, "08 / STEP 7")
h2(s, "トークンとコストを算出する — usage")
lead(s, "assistant レコードの message.usage に API と同じ粒度のトークン内訳が毎回記録されています。model と掛け合わせれば概算コストが出せます。")
codebox(s, 0.6, 2.2, 5.9, 2.6, [
    '"usage": {',
    [('  "input_tokens": 2,', {"color": CODE_K})],
    [('  "output_tokens": 410,', {"color": CODE_K})],
    '  "cache_creation_input_tokens": 39732,',
    '  "cache_read_input_tokens": 0,',
    '  "cache_creation": {',
    '    "ephemeral_1h_input_tokens": 39732,',
    '    "ephemeral_5m_input_tokens": 0 },',
    '  "service_tier": "standard" }',
], size=10.5)
codebox(s, 6.8, 2.2, 5.9, 2.6, [
    [("# セッション合計", {"color": CODE_C})],
    "jq -s '[.[] | select(.type==\"assistant\")",
    "        | .message.usage]",
    "  | { in:  (map(.input_tokens)|add),",
    "      out: (map(.output_tokens)|add),",
    "      cache_w: (map(.cache_creation_input_tokens)|add),",
    "      cache_r: (map(.cache_read_input_tokens)|add) }' \\",
    "  <session-id>.jsonl",
], size=10.5)
card(s, 0.6, 5.0, 12.1, 1.3, "掛け合わせ分析", [
    "モデル別（.message.model）に分けて単価を掛ければコスト概算／cache_read ÷ (input+cache_read) でキャッシュヒット率",
    "ブランチ・プロジェクト別に集計すれば「この機能開発に幾らかかったか」まで分かる。近道: npx ccusage が日別・モデル別レポートを自動生成",
])

# ---------------------------------------------------------------- 09 STEP8
s = base_slide("CC-LOG / 09", "STEP 8 — TIME")
eyebrow(s, "09 / STEP 8")
h2(s, "時間を分析する — timestamp の差分")
lead(s, "全レコードにミリ秒精度のタイムスタンプがあるので、差分を取るだけで応答速度・作業時間・活動パターンが数値化できます。")
card(s, 0.6, 2.15, 3.9, 1.75, "応答レイテンシ", [
    "!ユーザープロンプト → 直後の assistant レコードの時刻差。ターンごとの「待ち時間」。",
])
card(s, 4.7, 2.15, 3.9, 1.75, "ツール実行時間", [
    "!tool_use → 対応する tool_result の時刻差。重いビルドやテストが特定できる。",
])
card(s, 8.8, 2.15, 3.9, 1.75, "セッション時間", [
    "!最初と最後のレコードの差。日別・時間帯別に並べれば利用ヒートマップに。",
])
codebox(s, 0.6, 4.15, 12.1, 2.35, [
    [("import", {"color": CODE_K}), (" json, pandas ", {}), ("as", {"color": CODE_K}), (" pd", {})],
    [("rows = [json.loads(l) ", {}), ("for", {"color": CODE_K}), (" l ", {}), ("in", {"color": CODE_K}), (" open(path)]", {})],
    "df = pd.DataFrame([r for r in rows if 'timestamp' in r])",
    "df['ts'] = pd.to_datetime(df['timestamp'])",
    [("df['gap'] = df['ts'].diff()          ", {}), ("# レコード間隔", {"color": CODE_C})],
    [("df.groupby(df['ts'].dt.hour).size()   ", {}), ("# 時間帯別アクティビティ", {"color": CODE_C})],
])

# ---------------------------------------------------------------- 10 STEP9
s = base_slide("CC-LOG / 10", "STEP 9 — ERRORS")
eyebrow(s, "10 / STEP 9")
h2(s, "エラーと中断を洗い出す")
lead(s, "失敗パターンの分析は改善効果が最も大きい部分。ログには失敗のシグナルが複数の場所に残ります。")
table(s, 0.6, 2.1, 12.1, [2.6, 4.0, 5.5], [
    ["シグナル", "場所", "意味"],
    ["is_error: true", "user レコードの tool_result ブロック", "ツール呼び出しがエラーを返した"],
    ["stderr が非空", "toolUseResult", "コマンドは動いたが警告 / エラー出力あり"],
    ["interrupted: true", "toolUseResult", "ユーザーが実行を中断した（信頼のシグナル）"],
    ["stop_reason", "assistant の message", "tool_use / end_turn など応答の終わり方"],
    ["権限拒否の文言", "tool_result の本文", "許可プロンプトで拒否された操作の傾向"],
], row_h=0.42)
codebox(s, 0.6, 4.85, 12.1, 1.0, [
    [("# エラーになったツール呼び出しの数", {"color": CODE_C})],
    'jq \'[select(.type=="user") | .message.content[]? | select(.type=="tool_result" and .is_error==true)] | length\' *.jsonl',
])
memory(s, 0.6, 6.1, 12.1, 0.75, "使い方：",
       "「エラー率の高いツール×コマンド」を出すと、CLAUDE.md に書くべき注意書きが自動的に決まります。")

# ---------------------------------------------------------------- 11 STEP10
s = base_slide("CC-LOG / 11", "STEP 10 — VISUALIZE")
eyebrow(s, "11 / STEP 10")
h2(s, "可視化して1枚のレポートにする")
lead(s, "ここまでの集計を pandas に集めれば、あとは好きな形で出力するだけ。まずは4枚のチャートから始めるのがおすすめです。")
card(s, 0.6, 2.2, 5.5, 2.6, "最初に作る4チャート", [
    "日別トークン消費（積み上げ：入力 / 出力 / キャッシュ）",
    "ツール使用回数 Top10（横棒）",
    "時間帯×曜日の利用ヒートマップ",
    "セッションごとの往復数とエラー数（散布図）",
])
codebox(s, 6.3, 2.2, 6.4, 2.6, [
    [("import", {"color": CODE_K}), (" glob, json, pandas ", {}), ("as", {"color": CODE_K}), (" pd", {})],
    "recs = []",
    "for p in glob.glob('~/.claude/projects/*/*.jsonl'):",
    "    for l in open(p):",
    "        d = json.loads(l)",
    "        if d.get('type')=='assistant':",
    "            u = d['message'].get('usage',{})",
    "            recs.append({'ts':d['timestamp'], …})",
    [("df = pd.DataFrame(recs)  ", {}), ("# → groupby して plot", {"color": CODE_C})],
], size=10.5)
memory(s, 0.6, 5.15, 12.1, 1.0, "発展：",
       "この解析自体を Claude Code に頼むのが最速です。「~/.claude/projects の JSONL を読んで利用レポートを HTML で作って」— このスライドの手順がそのままプロンプトになります。")

# ---------------------------------------------------------------- 12 MATRIX
s = base_slide("CC-LOG / 12", "WHAT YOU CAN ANALYZE")
eyebrow(s, "12 / 全体マップ")
h2(s, "解析できること 一覧マトリクス")
table(s, 0.6, 1.75, 12.1, [1.6, 6.2, 4.3], [
    ["カテゴリ", "分かること", "使うフィールド"],
    ["会話", "プロンプト内容・往復数・依頼の種類・会話の分岐", "type=user, message.content, parentUuid"],
    ["行動", "ツール比率・実行コマンド・編集ファイル・委譲（サブエージェント）", "tool_use.name / input, isSidechain"],
    ["コスト", "トークン内訳・モデル別コスト・キャッシュヒット率", "usage.*, model, service_tier"],
    ["時間", "応答待ち・ツール実行時間・セッション長・活動時間帯", "timestamp の差分"],
    ["品質", "エラー率・中断率・権限拒否・stop_reason 分布", "is_error, stderr, interrupted, stop_reason"],
    ["文脈", "プロジェクト / ブランチ別の作業量・CC バージョン推移", "cwd, gitBranch, version"],
    ["思考", "thinking ブロックの量 ＝ 難易度の代理指標", 'content[].type=="thinking"'],
    ["タスク", "TodoWrite の計画と完了状況", "~/.claude/todos/*.json"],
], row_h=0.56)

# ---------------------------------------------------------------- 13 SNIPPETS
s = base_slide("CC-LOG / 13", "CHEAT SHEET")
eyebrow(s, "13 / チートシート")
h2(s, "コピペで使えるワンライナー集")
codebox(s, 0.6, 1.8, 5.95, 4.3, [
    [("# ① レコード種別の分布", {"color": CODE_C})],
    "jq -r '.type' s.jsonl | sort | uniq -c",
    "",
    [("# ② 人間のプロンプト一覧", {"color": CODE_C})],
    "jq -r 'select(.type==\"user\" and",
    "  (.message.content|type)==\"string\")",
    "  | .message.content' s.jsonl",
    "",
    [("# ③ ツール使用 Top10", {"color": CODE_C})],
    "jq -r 'select(.type==\"assistant\")",
    "  | .message.content[]?",
    "  | select(.type==\"tool_use\").name' s.jsonl \\",
    "  | sort | uniq -c | sort -rn | head",
], size=10.5)
codebox(s, 6.75, 1.8, 5.95, 4.3, [
    [("# ④ 総出力トークン", {"color": CODE_C})],
    "jq -s '[.[] | select(.type==\"assistant\")",
    "  | .message.usage.output_tokens] | add' s.jsonl",
    "",
    [("# ⑤ 使用モデルの内訳", {"color": CODE_C})],
    "jq -r 'select(.type==\"assistant\")",
    "  | .message.model' s.jsonl | sort | uniq -c",
    "",
    [("# ⑥ エラーになった tool_result 数", {"color": CODE_C})],
    "jq '[select(.type==\"user\")",
    "  | .message.content[]?",
    "  | select(.type==\"tool_result\"",
    "     and .is_error==true)] | length' s.jsonl",
], size=10.5)
memory(s, 0.6, 6.3, 12.1, 0.6, "メモ：",
       "s.jsonl は対象セッションファイル。*.jsonl にすればプロジェクト全体を一括集計できます。")

# ---------------------------------------------------------------- 14 WRAP
s = base_slide("CC-LOG / 14", "YOU CAN ANALYZE IT ALL")
eyebrow(s, "14 / 仕上げ")
h2(s, "まとめと、扱う上での注意")
card(s, 0.6, 2.0, 5.9, 3.6, "今日の手順（再掲）", [
    "在り処：~/.claude/projects/…/*.jsonl",
    "type で分類 → メタデータを軸にする",
    "parentUuid で会話ツリーを復元",
    "プロンプト / ツール / usage / 時刻を集計",
    "エラー・中断で改善点を特定",
    "pandas で1枚のレポートに可視化",
])
card(s, 6.8, 2.0, 5.9, 3.6, "注意点", [
    "!機密が入っている：コード片・パス・コマンド出力がそのまま残る。共有前に必ず内容を確認。",
    "!消える：cleanupPeriodDays（既定30日）で自動削除。分析用アーカイブは別途退避。",
    "!スキーマは非公式：バージョンで変わり得る。スクリプトは .get() で防御的に。",
])
memory(s, 0.6, 5.9, 12.1, 0.85, "次の一歩：",
       "週次で自動集計する cron / Routine を組めば「先週の Claude Code 利用レポート」が毎週届く仕組みまであと一歩です。")

out = "claude-code-log-analysis-slides.pptx"
prs.save(out)
print("saved", out)
